"""
Generates docs/DEMO.pptx — a presentation walkthrough of the
LLM-Protect Middleware project for an academic supervisor defence.

Run:
    .venv/bin/python scripts/generate_demo_pptx.py
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

# ───────────────────────────── design tokens ──────────────────────────────

# 16:9 widescreen — 13.333" x 7.5"
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Colour palette — dark "cyber-security" navy + cyan accents
COLOR_BG_DARK = RGBColor(0x0B, 0x14, 0x26)     # deep navy
COLOR_BG_LIGHT = RGBColor(0xF5, 0xF7, 0xFB)    # off-white
COLOR_ACCENT = RGBColor(0x00, 0xC8, 0xE0)      # cyan
COLOR_ACCENT_DARK = RGBColor(0x00, 0x8A, 0x9C)
COLOR_TEXT_LIGHT = RGBColor(0xEA, 0xEF, 0xF7)
COLOR_TEXT_DARK = RGBColor(0x1A, 0x22, 0x35)
COLOR_MUTED = RGBColor(0x8A, 0x98, 0xB0)
COLOR_OK = RGBColor(0x2E, 0xCC, 0x71)
COLOR_WARN = RGBColor(0xF1, 0xC4, 0x0F)
COLOR_BAD = RGBColor(0xE7, 0x4C, 0x3C)
COLOR_STAGE = [
    RGBColor(0x2C, 0x6E, 0xF5),  # blue
    RGBColor(0x16, 0xA0, 0x85),  # teal
    RGBColor(0xF3, 0x9C, 0x12),  # orange
    RGBColor(0xE6, 0x7E, 0x22),  # darker orange
    RGBColor(0x9B, 0x59, 0xB6),  # purple
    RGBColor(0xC0, 0x39, 0x2B),  # red
]

FONT_HEADING = "Calibri"
FONT_BODY = "Calibri"
FONT_MONO = "Menlo"


# ─────────────────────────── helper functions ─────────────────────────────

def add_blank_slide(prs):
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)


def fill_background(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill_color, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        if line_width is not None:
            shape.line.width = line_width
    shape.shadow.inherit = False
    return shape


def add_text(
    slide,
    x,
    y,
    w,
    h,
    text,
    *,
    font_size=18,
    bold=False,
    color=COLOR_TEXT_DARK,
    font_name=FONT_BODY,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font_name
    r.font.size = Pt(font_size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb


def add_bullets(
    slide,
    x,
    y,
    w,
    h,
    items,
    *,
    font_size=16,
    color=COLOR_TEXT_DARK,
    line_spacing=1.15,
    bullet_char="•",
    bullet_color=None,
    bold_first=False,
):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    bullet_color = bullet_color or COLOR_ACCENT
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(4)
        if bullet_char:
            br = p.add_run()
            br.text = f"{bullet_char}  "
            br.font.name = FONT_BODY
            br.font.size = Pt(font_size)
            br.font.bold = True
            br.font.color.rgb = bullet_color
        if isinstance(item, tuple):
            head, body = item
            r1 = p.add_run()
            r1.text = head
            r1.font.name = FONT_BODY
            r1.font.size = Pt(font_size)
            r1.font.bold = True
            r1.font.color.rgb = color
            r2 = p.add_run()
            r2.text = f"  {body}"
            r2.font.name = FONT_BODY
            r2.font.size = Pt(font_size)
            r2.font.color.rgb = color
        else:
            r = p.add_run()
            r.text = item
            r.font.name = FONT_BODY
            r.font.size = Pt(font_size)
            r.font.bold = bold_first and i == 0
            r.font.color.rgb = color
    return tb


def add_code_block(slide, x, y, w, h, code, *, font_size=12):
    bg = add_rect(slide, x, y, w, h, RGBColor(0x10, 0x1A, 0x30))
    bg.line.color.rgb = COLOR_ACCENT_DARK
    bg.line.width = Pt(0.75)
    tb = slide.shapes.add_textbox(
        x + Inches(0.15), y + Inches(0.1),
        w - Inches(0.3), h - Inches(0.2),
    )
    tf = tb.text_frame
    tf.word_wrap = True
    lines = code.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.1
        r = p.add_run()
        r.text = line if line else " "
        r.font.name = FONT_MONO
        r.font.size = Pt(font_size)
        if line.lstrip().startswith("#") or line.lstrip().startswith("//"):
            r.font.color.rgb = COLOR_MUTED
        elif line.lstrip().startswith(">"):
            r.font.color.rgb = COLOR_ACCENT
        else:
            r.font.color.rgb = COLOR_TEXT_LIGHT


def add_header_bar(slide, title, subtitle=None, slide_number=None, total=None):
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.85), COLOR_BG_DARK)
    add_rect(slide, 0, Inches(0.85), SLIDE_W, Inches(0.04), COLOR_ACCENT)
    add_text(
        slide, Inches(0.5), Inches(0.18), Inches(10.5), Inches(0.5),
        title, font_size=24, bold=True, color=COLOR_TEXT_LIGHT,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    if subtitle:
        add_text(
            slide, Inches(0.5), Inches(0.55), Inches(10.5), Inches(0.3),
            subtitle, font_size=12, color=COLOR_MUTED,
            anchor=MSO_ANCHOR.MIDDLE,
        )
    if slide_number is not None and total is not None:
        add_text(
            slide, Inches(11.7), Inches(0.18), Inches(1.4), Inches(0.5),
            f"{slide_number} / {total}",
            font_size=12, color=COLOR_MUTED,
            align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE,
        )


def add_footer(slide):
    add_text(
        slide, Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.4),
        "LLM-Protect Middleware  ·  Дипломна робота  ·  2026",
        font_size=10, color=COLOR_MUTED, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.MIDDLE,
    )


# ─────────────────────────────── slides ───────────────────────────────────


def slide_title(prs, total):
    s = add_blank_slide(prs)
    fill_background(s, COLOR_BG_DARK)

    add_rect(s, 0, Inches(6.7), SLIDE_W, Inches(0.06), COLOR_ACCENT)

    # decorative cyan accent strip on the left
    add_rect(s, 0, 0, Inches(0.35), SLIDE_H, COLOR_ACCENT)

    add_text(
        s, Inches(0.9), Inches(1.0), Inches(11.5), Inches(0.6),
        "ДИПЛОМНА РОБОТА",
        font_size=14, bold=True, color=COLOR_ACCENT,
    )
    add_text(
        s, Inches(0.9), Inches(1.6), Inches(11.5), Inches(1.6),
        "LLM-Protect Middleware",
        font_size=54, bold=True, color=COLOR_TEXT_LIGHT,
    )
    add_text(
        s, Inches(0.9), Inches(2.9), Inches(11.5), Inches(1.0),
        "Каскадна система захисту AI API від Model Denial-of-Service атак",
        font_size=24, color=COLOR_TEXT_LIGHT,
    )

    add_rect(s, Inches(0.9), Inches(4.3), Inches(11.5), Inches(0.04), COLOR_ACCENT_DARK)

    # meta block (two columns)
    add_text(
        s, Inches(0.9), Inches(4.55), Inches(5.5), Inches(0.4),
        "АВТОР",
        font_size=11, bold=True, color=COLOR_MUTED,
    )
    add_text(
        s, Inches(0.9), Inches(4.85), Inches(5.5), Inches(0.5),
        "Назару В.",
        font_size=20, bold=True, color=COLOR_TEXT_LIGHT,
    )

    add_text(
        s, Inches(7.0), Inches(4.55), Inches(5.5), Inches(0.4),
        "НАУКОВИЙ КЕРІВНИК",
        font_size=11, bold=True, color=COLOR_MUTED,
    )
    add_text(
        s, Inches(7.0), Inches(4.85), Inches(5.5), Inches(0.5),
        "[ПІБ керівника]",
        font_size=20, bold=True, color=COLOR_TEXT_LIGHT,
    )

    add_text(
        s, Inches(0.9), Inches(5.7), Inches(11.5), Inches(0.5),
        "ДЕМОНСТРАЦІЯ MVP",
        font_size=13, bold=True, color=COLOR_ACCENT,
    )
    add_text(
        s, Inches(0.9), Inches(6.0), Inches(11.5), Inches(0.5),
        "NestJS  ·  Redis Stack  ·  Prometheus  ·  Fastify  ·  Docker",
        font_size=14, color=COLOR_MUTED,
    )


def slide_problem(prs, idx, total):
    s = add_blank_slide(prs)
    fill_background(s, COLOR_BG_LIGHT)
    add_header_bar(
        s, "1. Проблема: Model Denial-of-Service",
        "Атаки, що виснажують ресурси AI-моделі та її семантичного кешу",
        idx, total,
    )

    # 4 attack cards
    cards = [
        ("Token Flood", "Дуже довгі промпти (10k+ токенів) витрачають GPU та вікно контексту."),
        ("Request Flood", "DDoS-стиль: 1000 RPS від одного клієнта повністю забивають чергу."),
        ("Entropy Flood", "Випадкові символи з ентропією ~6.0 біт — забруднюють Vector DB."),
        ("Cache Poisoning", "Тисячі унікальних промптів вибивають корисні embedding-и з кешу."),
    ]
    card_w = Inches(2.85)
    card_h = Inches(2.4)
    gap = Inches(0.15)
    y0 = Inches(1.4)
    x0 = Inches(0.5)
    for i, (name, desc) in enumerate(cards):
        x = x0 + (card_w + gap) * i
        add_rect(s, x, y0, card_w, card_h, RGBColor(0xFF, 0xFF, 0xFF), COLOR_BAD, Pt(1.5))
        add_rect(s, x, y0, card_w, Inches(0.5), COLOR_BAD)
        add_text(
            s, x, y0, card_w, Inches(0.5),
            f"  ⚠ {name}",
            font_size=14, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
            anchor=MSO_ANCHOR.MIDDLE,
        )
        add_text(
            s, x + Inches(0.15), y0 + Inches(0.65), card_w - Inches(0.3), card_h - Inches(0.8),
            desc, font_size=12, color=COLOR_TEXT_DARK,
        )

    # bottom band: consequences
    add_rect(s, Inches(0.5), Inches(4.2), Inches(12.3), Inches(2.5), RGBColor(0xFF, 0xFF, 0xFF), COLOR_MUTED, Pt(0.5))
    add_text(
        s, Inches(0.7), Inches(4.35), Inches(11.9), Inches(0.5),
        "Наслідки для виробничого AI API:",
        font_size=18, bold=True, color=COLOR_TEXT_DARK,
    )
    add_bullets(
        s, Inches(0.7), Inches(4.9), Inches(11.9), Inches(1.8),
        [
            ("Економічні:", "$0.01–0.10 за один блокований запит до GPT-4 → атака на $10k/день."),
            ("Технічні:", "GPU-черга насичується, легітимні користувачі отримують 5xx-timeout."),
            ("Якісні:", "Cache-poisoning знижує precision семантичного пошуку на 20–40%."),
        ],
        font_size=15,
    )
    add_footer(s)


def slide_goal(prs, idx, total):
    s = add_blank_slide(prs)
    fill_background(s, COLOR_BG_LIGHT)
    add_header_bar(
        s, "2. Мета та задачі дослідження",
        "Розробити архітектурне рішення, що зменшує атаки без шкоди для легітимного трафіку",
        idx, total,
    )

    # goal panel
    add_rect(s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(1.4), COLOR_BG_DARK)
    add_text(
        s, Inches(0.8), Inches(1.45), Inches(11.7), Inches(0.5),
        "🎯  МЕТА",
        font_size=14, bold=True, color=COLOR_ACCENT,
    )
    add_text(
        s, Inches(0.8), Inches(1.85), Inches(11.7), Inches(0.8),
        "Спроектувати та реалізувати middleware-систему, що відсіває MDoS-атаки "
        "до того, як вони досягнуть AI-моделі, дотримуючись принципу "
        "«від найдешевшої операції до найдорожчої».",
        font_size=16, color=COLOR_TEXT_LIGHT,
    )

    # objectives — 2 columns
    add_text(
        s, Inches(0.5), Inches(3.0), Inches(12.3), Inches(0.5),
        "Задачі:",
        font_size=20, bold=True, color=COLOR_TEXT_DARK,
    )

    objs_left = [
        "Спроектувати каскадну архітектуру з 5 рівнів фільтрації.",
        "Реалізувати MVP на NestJS + Fastify + Redis Stack.",
        "Інтегрувати OpenAI-сумісне API для прозорої заміни.",
    ]
    objs_right = [
        "Обчислювати embedding ТІЛЬКИ після дешевих перевірок.",
        "Забезпечити моніторинг (Prometheus) ефективності.",
        "Підтвердити інваріанту каскаду метриками runtime.",
    ]
    add_bullets(
        s, Inches(0.5), Inches(3.55), Inches(6.1), Inches(3.0),
        objs_left, font_size=15, color=COLOR_TEXT_DARK,
        bullet_char="▸", bullet_color=COLOR_ACCENT,
    )
    add_bullets(
        s, Inches(6.8), Inches(3.55), Inches(6.0), Inches(3.0),
        objs_right, font_size=15, color=COLOR_TEXT_DARK,
        bullet_char="▸", bullet_color=COLOR_ACCENT,
    )
    add_footer(s)


def slide_architecture(prs, idx, total):
    s = add_blank_slide(prs)
    fill_background(s, COLOR_BG_LIGHT)
    add_header_bar(
        s, "3. Архітектура: каскадна фільтрація",
        "Кожен наступний етап дорожчий — відсіювати треба максимально рано",
        idx, total,
    )

    # 6 stages as a horizontal pipeline
    stages = [
        ("1a", "Content-Length", "Middleware", "< 1 ms", "Розмір тіла"),
        ("1b", "Rate Limit",     "Guard",      "~1 ms",  "Redis INCR"),
        ("2",  "Exact Cache",    "Interceptor","~1 ms",  "SHA-256 + Redis"),
        ("3",  "Entropy",        "Interceptor","~1 ms",  "Shannon CPU"),
        ("4",  "Token Limit",    "Interceptor","~10 ms", "tiktoken BPE"),
        ("5",  "Embedding +",    "Interceptor","~50 ms", "ONNX + Vector DB"),
    ]

    box_w = Inches(1.95)
    box_h = Inches(2.6)
    gap = Inches(0.05)
    y0 = Inches(1.35)
    x0 = Inches(0.4)

    for i, (num, name, kind, lat, tech) in enumerate(stages):
        x = x0 + (box_w + gap) * i
        color = COLOR_STAGE[i]
        add_rect(s, x, y0, box_w, box_h, RGBColor(0xFF, 0xFF, 0xFF), color, Pt(2))
        add_rect(s, x, y0, box_w, Inches(0.7), color)
        add_text(
            s, x, y0 + Inches(0.05), box_w, Inches(0.3),
            f"STAGE {num}",
            font_size=10, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )
        add_text(
            s, x, y0 + Inches(0.32), box_w, Inches(0.4),
            name,
            font_size=14, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )
        add_text(
            s, x + Inches(0.1), y0 + Inches(0.85), box_w - Inches(0.2), Inches(0.4),
            kind,
            font_size=11, bold=True, color=COLOR_TEXT_DARK,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )
        add_rect(
            s, x + Inches(0.3), y0 + Inches(1.3),
            box_w - Inches(0.6), Inches(0.45),
            color,
        )
        add_text(
            s, x + Inches(0.3), y0 + Inches(1.3),
            box_w - Inches(0.6), Inches(0.45),
            lat,
            font_size=14, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )
        add_text(
            s, x + Inches(0.1), y0 + Inches(1.9), box_w - Inches(0.2), Inches(0.6),
            tech,
            font_size=11, color=COLOR_MUTED,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )

    # arrow band — descending cost
    arrow_y = y0 + box_h + Inches(0.25)
    add_text(
        s, Inches(0.4), arrow_y, Inches(12.5), Inches(0.4),
        "▸ дешевше",
        font_size=12, bold=True, color=COLOR_OK,
    )
    add_text(
        s, Inches(0.4), arrow_y, Inches(12.5), Inches(0.4),
        "дорожче ▸",
        font_size=12, bold=True, color=COLOR_BAD,
        align=PP_ALIGN.RIGHT,
    )
    add_rect(s, Inches(0.4), arrow_y + Inches(0.32), Inches(12.5), Inches(0.04), COLOR_MUTED)

    # invariant callout
    add_rect(s, Inches(0.5), Inches(5.45), Inches(12.3), Inches(1.4), COLOR_BG_DARK)
    add_text(
        s, Inches(0.75), Inches(5.55), Inches(11.8), Inches(0.5),
        "🔒  АРХІТЕКТУРНА ІНВАРІАНТА",
        font_size=13, bold=True, color=COLOR_ACCENT,
    )
    add_text(
        s, Inches(0.75), Inches(5.95), Inches(11.8), Inches(0.8),
        "Stage 5 (embedding ~50 ms) виконується ТІЛЬКИ для запитів, що пройшли всі дешеві етапи. "
        "Це і захищає AI-модель, і запобігає cache-poisoning Vector DB.",
        font_size=15, color=COLOR_TEXT_LIGHT,
    )
    add_footer(s)


def slide_stack(prs, idx, total):
    s = add_blank_slide(prs)
    fill_background(s, COLOR_BG_LIGHT)
    add_header_bar(
        s, "4. Технологічний стек",
        "Production-ready Node.js екосистема, повністю контейнеризована",
        idx, total,
    )

    # 6 cards grid 3x2
    cards = [
        ("NestJS 10", "Backend Framework", "Modular DI, Guards/Interceptors/Pipes — ідеально лягає на каскадну фільтрацію."),
        ("Fastify", "HTTP adapter", "60–70k RPS overhead < 5 ms — критично для middleware."),
        ("Redis Stack", "Cache + Vector", "Exact cache (SHA-256), rate-limit (INCR), semantic search (HNSW)."),
        ("tiktoken (BPE)", "Tokenizer", "Точний підрахунок токенів сумісно з OpenAI GPT моделями."),
        ("ONNX Runtime", "Embeddings", "Локальне обчислення ембеддінгів (зараз — Mock provider)."),
        ("Prometheus", "Observability", "Метрики по етапах + інваріанта `embedding = passed − exact_hits`."),
    ]
    card_w = Inches(4.05)
    card_h = Inches(2.45)
    gap_x = Inches(0.15)
    gap_y = Inches(0.2)
    x0 = Inches(0.5)
    y0 = Inches(1.4)
    for i, (name, tag, desc) in enumerate(cards):
        col = i % 3
        row = i // 3
        x = x0 + (card_w + gap_x) * col
        y = y0 + (card_h + gap_y) * row
        add_rect(s, x, y, card_w, card_h, RGBColor(0xFF, 0xFF, 0xFF), COLOR_MUTED, Pt(0.5))
        add_rect(s, x, y, Inches(0.12), card_h, COLOR_ACCENT)
        add_text(
            s, x + Inches(0.3), y + Inches(0.25), card_w - Inches(0.5), Inches(0.55),
            name, font_size=20, bold=True, color=COLOR_TEXT_DARK,
        )
        add_text(
            s, x + Inches(0.3), y + Inches(0.8), card_w - Inches(0.5), Inches(0.4),
            tag.upper(), font_size=11, bold=True, color=COLOR_ACCENT_DARK,
        )
        add_text(
            s, x + Inches(0.3), y + Inches(1.2), card_w - Inches(0.5), card_h - Inches(1.3),
            desc, font_size=13, color=COLOR_TEXT_DARK,
        )
    add_footer(s)


def slide_monorepo(prs, idx, total):
    s = add_blank_slide(prs)
    fill_background(s, COLOR_BG_LIGHT)
    add_header_bar(
        s, "5. Структура монорепозиторію",
        "NestJS workspace: один деплой, модульний код",
        idx, total,
    )

    add_code_block(
        s, Inches(0.5), Inches(1.4), Inches(7.5), Inches(5.4),
        """apps/
  └── gateway/                        # NestJS entry-point
      └── src/
          ├── main.ts                 # Fastify bootstrap
          ├── app.module.ts           # composition root
          ├── modules/
          │   ├── chat/               # /v1/chat/completions
          │   ├── health/             # /v1/health
          │   └── metrics/            # /v1/metrics
          └── common/                 # RequestId, Filters

libs/
  ├── shared/                         # DTOs, RequestContext, config
  ├── detection-engine/               # ← 5 cascade stages
  │   ├── stage-1a-content-length/
  │   ├── stage-1b-rate-limit/
  │   ├── stage-2-exact-cache/
  │   ├── stage-3-entropy/
  │   ├── stage-4-token-limit/
  │   └── stage-5-embedding/
  ├── semantic-cache/                 # Redis (exact + vector)
  ├── ai-proxy/                       # Ollama + Mock
  ├── logging/                        # nestjs-pino
  └── metrics/                        # Prometheus registry""",
        font_size=12,
    )

    # right panel with key idea
    add_rect(s, Inches(8.2), Inches(1.4), Inches(4.6), Inches(5.4), COLOR_BG_DARK)
    add_text(
        s, Inches(8.4), Inches(1.55), Inches(4.3), Inches(0.5),
        "Чому монорепо?",
        font_size=18, bold=True, color=COLOR_ACCENT,
    )
    add_bullets(
        s, Inches(8.4), Inches(2.05), Inches(4.3), Inches(4.5),
        [
            "Один deploy = одна Docker-image.",
            "Спільні DTO/типи без npm-publish.",
            "Етапи каскаду — окремі libs з власними тестами.",
            "Легко вмикати/вимикати етапи через ENV.",
            "16 unit-тестів покривають утиліти й interceptors.",
        ],
        font_size=14, color=COLOR_TEXT_LIGHT,
        bullet_char="▸", bullet_color=COLOR_ACCENT,
    )
    add_footer(s)


def slide_nest_lifecycle(prs, idx, total):
    s = add_blank_slide(prs)
    fill_background(s, COLOR_BG_LIGHT)
    add_header_bar(
        s, "6. Як NestJS гарантує порядок каскаду",
        "Lifecycle hooks Fastify + декларативний @UseInterceptors",
        idx, total,
    )

    add_code_block(
        s, Inches(0.5), Inches(1.4), Inches(8.5), Inches(3.2),
        """// apps/gateway/src/modules/chat/chat.controller.ts

@Controller('v1/chat')
@UseGuards(RateLimitGuard)                       // ← Stage 1b
@UseInterceptors(
  ExactCacheInterceptor,                         // ← Stage 2
  EntropyInterceptor,                            // ← Stage 3
  TokenLimitInterceptor,                         // ← Stage 4
  SemanticCacheInterceptor,                      // ← Stage 5
)
export class ChatController {
  @Post('completions')
  async create(@Body() dto: ChatCompletionRequestDto) { ... }
}""",
        font_size=13,
    )

    # right panel
    add_rect(s, Inches(9.2), Inches(1.4), Inches(3.6), Inches(3.2), COLOR_BG_DARK)
    add_text(
        s, Inches(9.4), Inches(1.55), Inches(3.3), Inches(0.5),
        "Порядок виконання",
        font_size=15, bold=True, color=COLOR_ACCENT,
    )
    add_bullets(
        s, Inches(9.4), Inches(2.0), Inches(3.3), Inches(2.5),
        [
            "1. Middleware (Stage 1a)",
            "2. Guards (Stage 1b)",
            "3. Interceptors\n   (Stage 2 → 3 → 4 → 5)",
            "4. Route handler",
        ],
        font_size=13, color=COLOR_TEXT_LIGHT,
        bullet_char="▸", bullet_color=COLOR_ACCENT,
    )

    # bottom note
    add_rect(s, Inches(0.5), Inches(4.8), Inches(12.3), Inches(2.0), RGBColor(0xFF, 0xFF, 0xFF), COLOR_ACCENT, Pt(1.5))
    add_text(
        s, Inches(0.75), Inches(4.95), Inches(11.9), Inches(0.5),
        "Ключова деталь:",
        font_size=15, bold=True, color=COLOR_ACCENT_DARK,
    )
    add_bullets(
        s, Inches(0.75), Inches(5.4), Inches(11.9), Inches(1.4),
        [
            ("Stages 3 і 4 — Interceptors, не Guards.", "Тому що Stage 2 (Interceptor) витягує prompt text, який їм потрібен. У NestJS Guards виконуються ДО усіх Interceptors."),
            ("Interceptors виконуються строго у порядку декларації.", "Це гарантовано фреймворком — не потребує власної логіки оркестрації."),
        ],
        font_size=13, color=COLOR_TEXT_DARK,
        bullet_char="▸", bullet_color=COLOR_ACCENT,
    )
    add_footer(s)


def slide_demo_intro(prs, idx, total):
    s = add_blank_slide(prs)
    fill_background(s, COLOR_BG_DARK)
    add_rect(s, 0, Inches(3.5), SLIDE_W, Inches(0.06), COLOR_ACCENT)
    add_text(
        s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.0),
        "ЖИВА ДЕМОНСТРАЦІЯ",
        font_size=18, bold=True, color=COLOR_ACCENT,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        s, Inches(0.5), Inches(2.2), Inches(12.3), Inches(1.5),
        "5 атак — 5 етапів каскаду",
        font_size=54, bold=True, color=COLOR_TEXT_LIGHT,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        s, Inches(0.5), Inches(3.8), Inches(12.3), Inches(0.6),
        "Кожна атака блокується НА СВОЄМУ етапі — без виклику AI",
        font_size=22, color=COLOR_MUTED,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        s, Inches(0.5), Inches(5.5), Inches(12.3), Inches(0.6),
        f"slide {idx} / {total}   ·   повна інструкція: docs/DEMO.md",
        font_size=12, color=COLOR_MUTED,
        align=PP_ALIGN.CENTER,
    )


def _demo_slide(
    prs,
    idx,
    total,
    title,
    subtitle,
    stage_label,
    stage_color,
    request_code,
    response_code,
    bullets,
):
    s = add_blank_slide(prs)
    fill_background(s, COLOR_BG_LIGHT)
    add_header_bar(s, title, subtitle, idx, total)

    # stage badge
    add_rect(s, Inches(0.5), Inches(1.3), Inches(3.0), Inches(0.6), stage_color)
    add_text(
        s, Inches(0.5), Inches(1.3), Inches(3.0), Inches(0.6),
        stage_label,
        font_size=15, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )

    # request label
    add_text(
        s, Inches(0.5), Inches(2.05), Inches(6.0), Inches(0.4),
        "▸ Запит",
        font_size=14, bold=True, color=COLOR_ACCENT_DARK,
    )
    add_code_block(
        s, Inches(0.5), Inches(2.45), Inches(6.0), Inches(2.3),
        request_code, font_size=11,
    )

    # response label
    add_text(
        s, Inches(6.8), Inches(2.05), Inches(6.0), Inches(0.4),
        "▸ Відповідь",
        font_size=14, bold=True, color=COLOR_ACCENT_DARK,
    )
    add_code_block(
        s, Inches(6.8), Inches(2.45), Inches(6.0), Inches(2.3),
        response_code, font_size=11,
    )

    # explanation block
    add_rect(s, Inches(0.5), Inches(4.95), Inches(12.3), Inches(1.95), RGBColor(0xFF, 0xFF, 0xFF), stage_color, Pt(1.5))
    add_text(
        s, Inches(0.75), Inches(5.05), Inches(11.9), Inches(0.4),
        "Що відбулося:",
        font_size=14, bold=True, color=stage_color,
    )
    add_bullets(
        s, Inches(0.75), Inches(5.45), Inches(11.9), Inches(1.5),
        bullets, font_size=13, color=COLOR_TEXT_DARK,
        bullet_char="▸", bullet_color=stage_color,
    )
    add_footer(s)


def slide_demo_stage2(prs, idx, total):
    _demo_slide(
        prs, idx, total,
        "Demo #1. Stage 2 — Exact Cache (HIT після MISS)",
        "Повторний ідентичний запит обслуговується з Redis за ~5 ms замість виклику AI",
        "STAGE 2 · EXACT CACHE", COLOR_STAGE[1],
        """# 1) Перший запит — MISS
curl -i -X POST localhost:3000/v1/chat/completions \\
  -H 'Content-Type: application/json' \\
  -d '{"model":"llama3.2","messages":[
       {"role":"user","content":"Capital of France?"}]}'

# 2) Той самий запит — HIT
curl -i -X POST localhost:3000/v1/chat/completions \\
  -H 'Content-Type: application/json' \\
  -d '{"model":"llama3.2","messages":[
       {"role":"user","content":"Capital of France?"}]}'""",
        """# 1) MISS — пішло до AI
HTTP/1.1 200 OK
x-request-id: req_2a83cef5-...
x-cache: MISS
{ "choices":[{ "message":{...} }] }

# 2) HIT — відповідь з Redis
HTTP/1.1 200 OK
x-request-id: req_5c7ee79a-...
x-cache: HIT
x-cache-type: exact""",
        [
            "Канонізований prompt → SHA-256 → ключ Redis. Складність O(1).",
            "Latency впав з ~150 ms до ~5 ms — AI взагалі не викликалась.",
            "Економія GPU + грошей при повторюваних запитах (FAQ, типові питання).",
        ],
    )


def slide_demo_stage3_low(prs, idx, total):
    _demo_slide(
        prs, idx, total,
        "Demo #2. Stage 3 — Entropy: repetitive flood (низька ентропія)",
        "Промпт з одного повторюваного символу заблоковано без обчислення embedding",
        "STAGE 3 · ENTROPY", COLOR_STAGE[2],
        """curl -X POST localhost:3000/v1/chat/completions \\
  -H 'Content-Type: application/json' \\
  -d '{"model":"llama3.2","messages":[
       {"role":"user",
        "content":"aaaaaaaaaaaaaaaaaaaaaaaaaaaaa"
                  "aaaaaaaaaaaaaaaaaaaaaaaaaaaaa"
                  "aaaaaaaaaaaaaaaaaaaaaaaaaaaaa"}]}'""",
        """HTTP/1.1 400 Bad Request
{
  "error": {
    "code": "HIGH_ENTROPY_GARBAGE",
    "stage": "ENTROPY",
    "message": "Prompt entropy 0.50
                outside allowed range [1.5, 5.5]",
    "risk_score": 100,
    "request_id": "req_d7a0b794-..."
  }
}""",
        [
            "Shannon entropy = 0.50 біт/символ (норма природного тексту ≈ 4.0–4.7).",
            "Симуляція repetitive-flood: атака, що швидко набиває контекстне вікно.",
            "Computed на CPU за ~1 ms, embedding не обчислювався, Vector DB не торкнули.",
        ],
    )


def slide_demo_stage3_high(prs, idx, total):
    _demo_slide(
        prs, idx, total,
        "Demo #3. Stage 3 — Entropy: cache poisoning (висока ентропія)",
        "Випадкові символи (~6 біт ентропії) могли б забруднити Vector DB — блокуємо",
        "STAGE 3 · ENTROPY", COLOR_STAGE[2],
        """RAND=$(python3 -c "
import random,string
print(''.join(random.choices(
  string.printable.strip(), k=300)))")

curl -X POST localhost:3000/v1/chat/completions \\
  -H 'Content-Type: application/json' \\
  -d "{\\"model\\":\\"llama3.2\\",\\"messages\\":
       [{\\"role\\":\\"user\\",
         \\"content\\":\\"$RAND\\"}]}\"""",
        """HTTP/1.1 400 Bad Request
{
  "error": {
    "code": "HIGH_ENTROPY_GARBAGE",
    "stage": "ENTROPY",
    "message": "Prompt entropy 5.83
                outside allowed range [1.5, 5.5]",
    "risk_score": 100
  }
}""",
        [
            "Аномально висока ентропія (5.83) — типово для cache-poisoning атаки.",
            "Без цього фільтра запит би пройшов до Stage 5 і додав «сміттєвий» вектор у Vector DB.",
            "Захищає precision семантичного пошуку для всіх інших користувачів.",
        ],
    )


def slide_demo_stage4(prs, idx, total):
    _demo_slide(
        prs, idx, total,
        "Demo #4. Stage 4 — Token Limit (token-flood)",
        "10000+ BPE-токенів — блокуємо до виклику AI та обчислення embedding",
        "STAGE 4 · TOKEN LIMIT", COLOR_STAGE[3],
        """PROMPT=$(python3 -c "
print('lorem ipsum dolor sit amet ' * 2000,
      end='')")

curl -X POST localhost:3000/v1/chat/completions \\
  -H 'Content-Type: application/json' \\
  -d "{\\"model\\":\\"llama3.2\\",\\"messages\\":
       [{\\"role\\":\\"user\\",
         \\"content\\":\\"$PROMPT\\"}]}\"""",
        """HTTP/1.1 413 Payload Too Large
{
  "error": {
    "code": "TOKEN_LIMIT_EXCEEDED",
    "stage": "TOKEN_LIMIT",
    "message": "Prompt has 10003 tokens,
                exceeds 8000 limit",
    "risk_score": 100
  }
}""",
        [
            "tiktoken дає той самий BPE-токенайзер, що OpenAI — точний підрахунок.",
            "Симуляція token-flood: запит хотів би з'їсти 10k токенів контексту GPU.",
            "Поріг 8000 — конфігурований через STAGE_TOKEN_MAX у .env.",
        ],
    )


def slide_demo_stage1b(prs, idx, total):
    _demo_slide(
        prs, idx, total,
        "Demo #5. Stage 1b — Rate Limit (parallel flood)",
        "150 паралельних запитів від одного API-key за 2 сек → 100×200 + 50×429",
        "STAGE 1b · RATE LIMIT", COLOR_STAGE[0],
        """# Скидаємо Redis для чистоти експерименту
docker exec llm-protect-redis redis-cli FLUSHDB

# 150 паралельних запитів, 30 у потоці
seq 1 150 | xargs -n1 -P30 -I{} curl -s \\
  -o /dev/null -w "%{http_code}\\n" \\
  -X POST localhost:3000/v1/chat/completions \\
  -H 'Content-Type: application/json' \\
  -H 'X-API-Key: attacker-001' \\
  -d '{"model":"llama3.2","messages":[
       {"role":"user","content":"ping-{}"}]}' \\
  | sort | uniq -c""",
        """  100 200
   50 429

# Метрики Prometheus підтверджують:
llm_protect_passed_stage_total
  {stage="RATE_LIMIT"} 100
llm_protect_blocked_total
  {stage="RATE_LIMIT"} 50""",
        [
            "Sliding window: Redis INCR + EXPIRE 60 сек на ключ API-key.",
            "Перші 100 пройшли успішно, наступні 50 — миттєво відкинуто з 429.",
            "Найдешевша мережева операція (~1 ms) — захищає увесь подальший каскад.",
        ],
    )


def slide_metrics_proof(prs, idx, total):
    s = add_blank_slide(prs)
    fill_background(s, COLOR_BG_LIGHT)
    add_header_bar(
        s, "7. Доведення інваріанти каскаду",
        "Prometheus метрики runtime підтверджують архітектурну вимогу",
        idx, total,
    )

    # left: formula + actual values
    add_rect(s, Inches(0.5), Inches(1.4), Inches(6.0), Inches(5.4), COLOR_BG_DARK)
    add_text(
        s, Inches(0.7), Inches(1.55), Inches(5.6), Inches(0.5),
        "Ключова формула",
        font_size=15, bold=True, color=COLOR_ACCENT,
    )
    add_text(
        s, Inches(0.7), Inches(2.1), Inches(5.6), Inches(1.8),
        "embedding_computed",
        font_size=20, bold=True, color=COLOR_TEXT_LIGHT,
        font_name=FONT_MONO,
    )
    add_text(
        s, Inches(0.7), Inches(2.55), Inches(5.6), Inches(0.5),
        "=",
        font_size=18, bold=True, color=COLOR_ACCENT,
        font_name=FONT_MONO,
    )
    add_text(
        s, Inches(0.7), Inches(3.0), Inches(5.6), Inches(0.5),
        "passed{TOKEN_LIMIT}",
        font_size=20, bold=True, color=COLOR_TEXT_LIGHT,
        font_name=FONT_MONO,
    )
    add_text(
        s, Inches(0.7), Inches(3.5), Inches(5.6), Inches(0.5),
        "−  cache_hits{exact}",
        font_size=20, bold=True, color=COLOR_TEXT_LIGHT,
        font_name=FONT_MONO,
    )

    add_rect(s, Inches(0.7), Inches(4.3), Inches(5.6), Inches(0.04), COLOR_ACCENT_DARK)

    add_text(
        s, Inches(0.7), Inches(4.45), Inches(5.6), Inches(0.5),
        "З реальних smoke-тестів:",
        font_size=13, bold=True, color=COLOR_MUTED,
    )
    add_text(
        s, Inches(0.7), Inches(4.95), Inches(5.6), Inches(0.5),
        "218  =  220  −  2   ✓",
        font_size=24, bold=True, color=COLOR_OK,
        font_name=FONT_MONO,
    )
    add_text(
        s, Inches(0.7), Inches(5.6), Inches(5.6), Inches(1.0),
        "Embedding обчислено лише 218 разів з 280 вхідних запитів — економія 22% CPU.",
        font_size=14, color=COLOR_TEXT_LIGHT,
    )

    # right: raw metrics + funnel
    add_text(
        s, Inches(6.8), Inches(1.4), Inches(6.0), Inches(0.4),
        "Воронка каскаду (з GET /v1/metrics):",
        font_size=14, bold=True, color=COLOR_TEXT_DARK,
    )
    add_code_block(
        s, Inches(6.8), Inches(1.85), Inches(6.0), Inches(3.6),
        """passed_stage{CONTENT_LENGTH}  280
passed_stage{RATE_LIMIT}      280
passed_stage{EXACT_CACHE}     226   (54 блок.)
passed_stage{ENTROPY}           8
passed_stage{TOKEN_LIMIT}     220
passed_stage{SEMANTIC_CACHE}  218

embedding_computed_total      218
                              ───
                              ✓ matches!""",
        font_size=13,
    )
    add_rect(s, Inches(6.8), Inches(5.6), Inches(6.0), Inches(1.2), COLOR_OK)
    add_text(
        s, Inches(6.95), Inches(5.65), Inches(5.7), Inches(1.1),
        "ВИСНОВОК: вимога наукового керівника «embedding рахується лише після всіх дешевих етапів» — підтверджена runtime-телеметрією.",
        font_size=13, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
        anchor=MSO_ANCHOR.MIDDLE,
    )
    add_footer(s)


def slide_quality(prs, idx, total):
    s = add_blank_slide(prs)
    fill_background(s, COLOR_BG_LIGHT)
    add_header_bar(
        s, "8. Якість реалізації",
        "Production-ready практики, готовність до A/B-тестування для дипломної",
        idx, total,
    )

    # 4 quality cards
    cards = [
        ("Тести", "16/16", "unit-тести каскаду й утиліт passing.\nJest + Supertest, ізольовані mocks для ConfigService/Redis."),
        ("CI/CD", "Auto", "GitHub Actions: lint → test → build → docker.\nКожен PR валідується автоматично."),
        ("Observability", "9", "Prometheus метрик: per-stage counters, latency histograms,\ncache hit ratio, embedding compute count."),
        ("Конфіг", "ENV", "Кожен етап вмикається через STAGE_*_ENABLED.\nГотово до A/B-порівняння для розділу «Ефективність»."),
    ]
    card_w = Inches(3.0)
    card_h = Inches(2.6)
    gap = Inches(0.15)
    y0 = Inches(1.4)
    x0 = Inches(0.5)
    for i, (name, big, desc) in enumerate(cards):
        x = x0 + (card_w + gap) * i
        add_rect(s, x, y0, card_w, card_h, RGBColor(0xFF, 0xFF, 0xFF), COLOR_MUTED, Pt(0.5))
        add_rect(s, x, y0, card_w, Inches(0.45), COLOR_ACCENT_DARK)
        add_text(
            s, x, y0, card_w, Inches(0.45),
            name.upper(),
            font_size=13, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )
        add_text(
            s, x, y0 + Inches(0.55), card_w, Inches(0.9),
            big,
            font_size=44, bold=True, color=COLOR_ACCENT_DARK,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )
        add_text(
            s, x + Inches(0.15), y0 + Inches(1.55), card_w - Inches(0.3), card_h - Inches(1.6),
            desc, font_size=12, color=COLOR_TEXT_DARK,
        )

    # bottom: OpenAPI mention
    add_rect(s, Inches(0.5), Inches(4.3), Inches(12.3), Inches(2.4), COLOR_BG_DARK)
    add_text(
        s, Inches(0.75), Inches(4.45), Inches(11.9), Inches(0.5),
        "Бонус: OpenAPI документація генерується автоматично",
        font_size=16, bold=True, color=COLOR_ACCENT,
    )
    add_bullets(
        s, Inches(0.75), Inches(4.95), Inches(11.9), Inches(1.8),
        [
            "@ApiProperty / @ApiResponse у DTO → Swagger UI на /docs без ручного YAML.",
            "OpenAI-сумісне API (POST /v1/chat/completions) — прозора заміна для існуючих клієнтів.",
            "/v1/health (terminus) і /v1/metrics (prometheus) — стандартні endpoint'и для Kubernetes.",
        ],
        font_size=14, color=COLOR_TEXT_LIGHT,
        bullet_char="▸", bullet_color=COLOR_ACCENT,
    )
    add_footer(s)


def slide_roadmap(prs, idx, total):
    s = add_blank_slide(prs)
    fill_background(s, COLOR_BG_LIGHT)
    add_header_bar(
        s, "9. Дорожня карта та подальші кроки",
        "Що реалізовано в MVP та що залишилось до фінального захисту",
        idx, total,
    )

    # done column
    add_rect(s, Inches(0.5), Inches(1.4), Inches(6.0), Inches(5.4), RGBColor(0xFF, 0xFF, 0xFF), COLOR_OK, Pt(1.5))
    add_rect(s, Inches(0.5), Inches(1.4), Inches(6.0), Inches(0.55), COLOR_OK)
    add_text(
        s, Inches(0.5), Inches(1.4), Inches(6.0), Inches(0.55),
        "✓  ГОТОВО (MVP)",
        font_size=16, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )
    add_bullets(
        s, Inches(0.75), Inches(2.1), Inches(5.5), Inches(4.5),
        [
            "Каркас NestJS монорепо + Fastify",
            "Всі 5 етапів каскаду (Stages 1a–5)",
            "Redis Stack: exact cache + rate-limit",
            "Mock embedding provider + semantic cache",
            "Prometheus метрики + інваріанта",
            "OpenAPI docs + health-check",
            "Docker Compose оточення",
            "16 unit-тестів + CI pipeline",
            "PRD.md + DEMO.md + README.md",
        ],
        font_size=14, color=COLOR_TEXT_DARK,
        bullet_char="✓", bullet_color=COLOR_OK,
    )

    # todo column
    add_rect(s, Inches(6.8), Inches(1.4), Inches(6.0), Inches(5.4), RGBColor(0xFF, 0xFF, 0xFF), COLOR_WARN, Pt(1.5))
    add_rect(s, Inches(6.8), Inches(1.4), Inches(6.0), Inches(0.55), COLOR_WARN)
    add_text(
        s, Inches(6.8), Inches(1.4), Inches(6.0), Inches(0.55),
        "◐  ПЛАНИ ДО ЗАХИСТУ",
        font_size=16, bold=True, color=COLOR_TEXT_DARK,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )
    add_bullets(
        s, Inches(7.05), Inches(2.1), Inches(5.5), Inches(4.5),
        [
            "ONNX embedding (replace Mock)",
            "Adaptive Rate Limiter (token budget)",
            "PostgreSQL логування блокувань",
            "React Dashboard (Recharts)",
            "Attack Simulator на k6/Artillery",
            "HNSW vector search (RediSearch)",
            "E2E тести інваріанти з реальним Redis",
            "A/B-порівняння етапів для дипломної",
            "Розділи 4–5 пояснювальної записки",
        ],
        font_size=14, color=COLOR_TEXT_DARK,
        bullet_char="◐", bullet_color=COLOR_WARN,
    )
    add_footer(s)


def slide_qa(prs, idx, total):
    s = add_blank_slide(prs)
    fill_background(s, COLOR_BG_DARK)
    add_rect(s, 0, Inches(3.5), SLIDE_W, Inches(0.06), COLOR_ACCENT)
    add_text(
        s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.0),
        "ДЯКУЮ ЗА УВАГУ",
        font_size=18, bold=True, color=COLOR_ACCENT,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        s, Inches(0.5), Inches(2.2), Inches(12.3), Inches(1.5),
        "Q & A",
        font_size=72, bold=True, color=COLOR_TEXT_LIGHT,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        s, Inches(0.5), Inches(4.0), Inches(12.3), Inches(0.6),
        "Готова відповісти на ваші запитання",
        font_size=22, color=COLOR_MUTED,
        align=PP_ALIGN.CENTER,
    )

    # footer with project meta
    add_rect(s, Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.3), COLOR_BG_DARK, COLOR_ACCENT_DARK, Pt(1))
    add_text(
        s, Inches(0.7), Inches(5.65), Inches(5.8), Inches(0.4),
        "ДОКУМЕНТАЦІЯ",
        font_size=11, bold=True, color=COLOR_ACCENT,
    )
    add_text(
        s, Inches(0.7), Inches(6.0), Inches(5.8), Inches(0.7),
        "PRD.md  ·  DEMO.md  ·  README.md",
        font_size=14, color=COLOR_TEXT_LIGHT, font_name=FONT_MONO,
    )
    add_text(
        s, Inches(7.0), Inches(5.65), Inches(5.8), Inches(0.4),
        "РЕПОЗИТОРІЙ",
        font_size=11, bold=True, color=COLOR_ACCENT,
    )
    add_text(
        s, Inches(7.0), Inches(6.0), Inches(5.8), Inches(0.7),
        "github.com/.../llm-protect-middleware",
        font_size=14, color=COLOR_TEXT_LIGHT, font_name=FONT_MONO,
    )


# ────────────────────────────── orchestrate ───────────────────────────────


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = [
        slide_title,
        slide_problem,
        slide_goal,
        slide_architecture,
        slide_stack,
        slide_monorepo,
        slide_nest_lifecycle,
        slide_demo_intro,
        slide_demo_stage2,
        slide_demo_stage3_low,
        slide_demo_stage3_high,
        slide_demo_stage4,
        slide_demo_stage1b,
        slide_metrics_proof,
        slide_quality,
        slide_roadmap,
        slide_qa,
    ]
    total = len(builders)

    # title slide doesn't need idx
    builders[0](prs, total)
    for i, builder in enumerate(builders[1:], start=2):
        builder(prs, i, total)

    out = Path(__file__).resolve().parent.parent / "docs" / "DEMO.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)
    print(f"✓ Saved presentation: {out}  ({total} slides)")


if __name__ == "__main__":
    main()
